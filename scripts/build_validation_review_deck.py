from __future__ import annotations

import json
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
VALIDATION = ROOT / "docs/production_runs/full_all_intake_as_new_20260520_consolidated/report_validation.json"
CONFLICTS = ROOT / "data/calibration/report_baseline_conflicts.json"
OUTPUT = Path(r"C:\Users\duxy\Desktop\duxyb\CableTrayAI_历史报告冲突原因分析.pptx")


INK = RGBColor(23, 33, 43)
MUTED = RGBColor(100, 116, 139)
GREEN = RGBColor(15, 118, 110)
RED = RGBColor(180, 35, 24)
AMBER = RGBColor(183, 121, 31)
BLUE = RGBColor(29, 78, 216)
BG = RGBColor(248, 250, 252)
LINE = RGBColor(203, 213, 225)
DARK = RGBColor(15, 23, 42)


def read_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8-sig"))


def pct(value: Any) -> str:
    try:
        return f"{float(value) * 100:.2f}%"
    except (TypeError, ValueError):
        return "-"


def clean(value: Any, limit: int = 180) -> str:
    text = " ".join(str(value or "").split())
    return text if len(text) <= limit else text[: limit - 1] + "…"


def set_font(paragraph, *, size: float = 10, bold: bool = False, color: RGBColor = INK) -> None:
    for run in paragraph.runs:
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_text(slide, text: str, x: float, y: float, w: float, h: float, *, size: float = 10, bold: bool = False, color: RGBColor = INK, align: PP_ALIGN | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align is not None:
        paragraph.alignment = align
    set_font(paragraph, size=size, bold=bold, color=color)


def add_title(slide, title: str, kicker: str, page: int) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_text(slide, kicker, 0.55, 0.28, 3.2, 0.25, size=7.5, bold=True, color=GREEN)
    add_text(slide, title, 0.55, 0.58, 11.8, 0.45, size=21, bold=True, color=INK)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.14), Inches(12.1), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    add_text(slide, "CableTrayAI | 后验验证问题必须保留，不得伪装为通过", 0.55, 7.08, 9.5, 0.2, size=7.3, color=MUTED)
    add_text(slide, str(page).zfill(2), 12.0, 7.05, 0.7, 0.2, size=7.3, color=MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, x: float, y: float, w: float, h: float, title: str, body: str, *, fill: RGBColor = RGBColor(255, 255, 255), line: RGBColor = LINE, title_color: RGBColor = INK, body_size: float = 9.0) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    add_text(slide, title, x + 0.16, y + 0.12, w - 0.32, 0.24, size=10.5, bold=True, color=title_color)
    add_text(slide, body, x + 0.16, y + 0.46, w - 0.32, h - 0.55, size=body_size, color=INK)


def add_table(slide, rows: list[list[Any]], x: float, y: float, w: float, h: float, widths: list[float], *, font_size: float = 7.5) -> None:
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c_idx == 0 or len(str(value)) < 12 else PP_ALIGN.LEFT
                set_font(paragraph, size=font_size, bold=(r_idx == 0), color=INK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(241, 245, 249) if r_idx == 0 else RGBColor(255, 255, 255)


def add_bar_chart(slide, items: list[tuple[str, int]], x: float, y: float, w: float, h: float, color: RGBColor) -> None:
    max_value = max([value for _, value in items] or [1])
    row_h = h / max(len(items), 1)
    for index, (label, value) in enumerate(items):
        yy = y + index * row_h
        add_text(slide, label, x, yy, 2.6, row_h * 0.8, size=7.3, color=INK)
        bar_w = (w - 3.35) * value / max_value
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 2.75), Inches(yy + 0.04), Inches(bar_w), Inches(row_h * 0.55))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.color.rgb = color
        add_text(slide, str(value), x + w - 0.5, yy, 0.45, row_h * 0.8, size=7.2, color=MUTED, align=PP_ALIGN.RIGHT)


def summarize_failed_cases(failed_cases: list[dict[str, Any]]) -> list[list[str]]:
    rows = [["报告号", "最大误差", "主要超差类型"]]
    for item in failed_cases:
        counts = Counter(metric.get("metric_type", "unknown") for metric in item.get("failed_metrics", []))
        top = "; ".join(f"{key}:{value}" for key, value in counts.most_common(4))
        rows.append([item.get("report_no", ""), pct(item.get("max_gate_error")), clean(top, 72)])
    return rows


def main() -> None:
    validation = read_json(VALIDATION)
    conflicts = read_json(CONFLICTS).get("conflicts", [])
    results = validation.get("results", [])
    pass_cases = [item for item in results if item.get("status") == "pass"]
    conflict_cases = [item for item in results if item.get("status") == "baseline_conflict"]
    failed_cases = [item for item in results if item.get("status") == "fail"]
    failed_metrics = [metric for item in failed_cases for metric in item.get("failed_metrics", [])]
    failed_counts = Counter(metric.get("metric_type", "unknown") for metric in failed_metrics)
    conflict_package_counts = Counter(item.get("package_id", "unknown") for item in conflicts)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    page = 1

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK
    add_text(slide, "CableTrayAI", 0.65, 0.46, 2.2, 0.3, size=12, bold=True, color=RGBColor(167, 243, 208))
    add_text(slide, "历史报告冲突与生产计算逻辑审查", 0.65, 1.02, 10.5, 0.7, size=30, bold=True, color=RGBColor(255, 255, 255))
    add_text(slide, "15 个已审历史报告/源文件冲突与 24 个仍超 1% 后验验证问题", 0.68, 2.0, 10.8, 0.35, size=13, color=RGBColor(226, 232, 240))
    add_card(slide, 0.75, 3.05, 3.1, 1.35, "严格通过", f"{len(pass_cases)} 个\n与历史报告在 1% 门禁内一致", fill=RGBColor(236, 253, 243), line=RGBColor(141, 214, 191), title_color=GREEN, body_size=15)
    add_card(slide, 4.05, 3.05, 3.2, 1.35, "已审冲突", f"{len(conflict_cases)} 个\n有证据说明报告/源文件互相冲突", fill=RGBColor(255, 250, 235), line=RGBColor(242, 201, 76), title_color=AMBER, body_size=13)
    add_card(slide, 7.45, 3.05, 3.2, 1.35, "仍需审查", f"{len(failed_cases)} 个\n不能假装通过，保留为后验验证问题", fill=RGBColor(254, 243, 242), line=RGBColor(253, 162, 155), title_color=RED, body_size=13)
    add_text(slide, "原则：报告对比用于校准与发现冲突；生产结果以提资事实、标准命令流、谱文件、真实 ANSYS 输出和确定性评定为准。", 0.75, 5.42, 11.6, 0.35, size=12, color=RGBColor(226, 232, 240))
    page += 1

    slide = prs.slides.add_slide(blank)
    add_title(slide, "47 个历史报告样例的当前状态", "VALIDATION STATUS", page)
    add_bar_chart(slide, [("严格通过", len(pass_cases)), ("已审冲突", len(conflict_cases)), ("仍超 1%", len(failed_cases))], 0.8, 1.6, 5.8, 1.9, GREEN)
    add_table(
        slide,
        [
            ["分类", "数量", "处理口径"],
            ["严格通过", len(pass_cases), "可作为正向回归样例，但仍保留 source_ref 和 result_validation。"],
            ["已审冲突", len(conflict_cases), "有证据说明历史报告、源命令流、备份源或提资事实无法同时成立；不用于硬校准。"],
            ["仍超 1%", len(failed_cases), "后验验证问题；继续按模型拓扑、谱系数、节点集、材料/公式逐项审。"],
        ],
        6.9,
        1.45,
        5.75,
        2.25,
        [1.2, 0.7, 3.85],
        font_size=8.0,
    )
    add_card(slide, 0.8, 4.35, 11.85, 1.35, "口径边界", "15 个已审冲突不是通过；24 个未过也不是简单软件失败。两类问题都必须保留，但不能用历史报告数字反向硬改生产命令流。", body_size=11)
    page += 1

    slide = prs.slides.add_slide(blank)
    add_title(slide, "生产计算链路：哪些代码负责什么", "CODE RESPONSIBILITY MAP", page)
    add_table(
        slide,
        [
            ["模块", "职责", "关键风险控制"],
            ["core/intake", "解析提资 Excel，抽取工程号、厂房、标高、托盘层数、荷载、方钢截面。", "不按固定行号；报告号缺失时用临时提资身份。"],
            ["core/spectra", "读取反应谱工作簿，按厂房/标高/阻尼/SL-1/SL-2/方向选择或插值。", "谱配置未确认不得 real-run。"],
            ["core/apdl", "由标准 APDL/PIP/MAC 源族生成三份命令流：建模、计算、提取。", "只做参数化和必要拓扑对齐，不靠近似值挑结果。"],
            ["core/ansys", "自动发现 ANSYS，preflight，run_all.mac，真实运行和审计。", "失败不得切 mock；real-run 必须显式确认。"],
            ["core/results", "解析 LIS/OUP/BMP/PNG，生成 result.json 和 figures_manifest.json。", "全零、UNKNOWN 节点、缺图、空白图一律阻断。"],
            ["core/evaluators", "确定性公式和 Excel 权威评定。", "公式必须有 source_ref；未确认不得最终通过。"],
            ["core/report", "把 result.json 和图片注入 Word 模板。", "不改正文结构和模板格式，只替换表格数据与图片。"],
        ],
        0.55,
        1.35,
        12.25,
        5.45,
        [1.55, 5.05, 5.65],
        font_size=6.8,
    )
    page += 1

    slide = prs.slides.add_slide(blank)
    add_title(slide, "已审冲突不是“调参通过”：它们都有证据链", "AUDITED CONFLICTS", page)
    add_bar_chart(slide, list(conflict_package_counts.items()), 0.8, 1.45, 4.3, 2.05, AMBER)
    add_table(
        slide,
        [
            ["样例", "冲突性质", "结论"],
            ["4123", "谱/静力系数冲突：源和谱用竖向约 0.409，报告部分行像用水平约 0.378。", "生产保留谱文件系数，不为贴报告压平竖向。"],
            ["4126", "紧凑源缺少标准 *9 连接螺栓代理节点族；补入会改变模态/支架应力。", "仅把连接螺栓载荷行列为历史冲突。"],
            ["4128", "两行连接载荷完整复制 4115；另有 FZ 像来自 .bak 备用源。", "仅排除可疑历史基准行，不反向改生产结果。"],
        ],
        5.4,
        1.35,
        7.25,
        3.0,
        [0.85, 3.9, 2.5],
        font_size=7.5,
    )
    add_card(slide, 0.8, 4.65, 11.85, 0.95, "审查要求", "这些冲突仍需科室签字确认：历史报告是否接受、是否需要修订、是否作为 legacy conflict 保留。软件侧不把这些冲突用于新提资生产调参。", fill=RGBColor(255, 247, 237), line=RGBColor(253, 186, 116), body_size=10.3)
    page += 1

    for title, kicker, rows, note in [
        (
            "4123：谱系数冲突导致多表联动偏差",
            "CASE 4123",
            [
                ["不一致位置", "支架/托臂应力、焊缝等效、基础载荷、连接载荷的 SL-1 相关行。"],
                ["证据", "active solve 与谱文件使用竖向系数约 0.409；报告结果更接近水平系数约 0.378。"],
                ["处理", "登记为历史源/报告冲突。生产计算继续使用所选谱文件，不做历史报告反向修正。"],
            ],
            "如果科室确认报告应按竖向谱系数，则历史报告需要标记或修订；如果确认历史命令流另有依据，需要补充该依据进入 source_ref。",
        ),
        (
            "4126：连接螺栓载荷节点集与模型拓扑冲突",
            "CASE 4126",
            [
                ["紧凑历史源", "能复现模态、方钢应力、基础载荷、托臂、焊缝，但缺少标准 *9 托盘-托臂连接代理节点族。"],
                ["强行补标准 *9 节点", "连接载荷提取路径更完整，但会改变模态和支架应力，说明不是同一个可审计模型。"],
                ["生产逻辑", "按标准拓扑生成并输出 result_source_map；全零、UNKNOWN、缺节点直接 fail，不贴报告选近值。"],
            ],
            "4126 的问题不是公式猜错，而是历史源模型和报告连接载荷集合无法用同一个拓扑同时解释。",
        ),
        (
            "4128：报告行疑似复制和备用源混用",
            "CASE 4128",
            [
                ["连接螺栓载荷表两行", "全行数值与 4115 报告对应行完全重复。"],
                ["派生螺栓 FZ", "FX/FY/MY/MZ 与 active source 接近，只有 FZ 更像 .bak 源族。"],
                ["其他表", "源生成的 HF-FORCE、JCZH、应力、模态与报告大体一致，只排除冲突行。"],
            ],
            "这类冲突不允许通过“哪个数接近就用哪个”解决。必须登记 source_ref、冲突证据和科室确认意见。",
        ),
    ]:
        slide = prs.slides.add_slide(blank)
        add_title(slide, title, kicker, page)
        add_table(slide, [["项目", "说明"], *rows], 0.8, 1.45, 11.8, 3.1, [2.4, 9.4], font_size=8.2)
        add_card(slide, 0.8, 5.05, 11.8, 0.9, "审查结论", note, fill=RGBColor(255, 255, 255), body_size=10.8)
        page += 1

    slide = prs.slides.add_slide(blank)
    add_title(slide, "24 个仍超 1%：失败指标集中在哪里", "REMAINING MISMATCHES", page)
    add_bar_chart(slide, failed_counts.most_common(9), 0.75, 1.35, 6.25, 4.65, RED)
    add_table(
        slide,
        [
            ["优先级", "看什么", "为什么"],
            ["1", "modal_frequency 大偏差", "模态大偏差通常说明模型拓扑/约束/质量源不同，不应先调公式。"],
            ["2", "beam_calculation_value", "梁应力是组合比和焊缝等效的上游，先审元素集合、截面、材料、谱系数。"],
            ["3", "foundation / connection load", "审 JCZH/LS-FORCE 的节点集合和 envelope 规则。"],
            ["4", "evaluation_ratio", "审许用值、事故系数、Excel 单元格来源和四舍五入。"],
        ],
        7.2,
        1.55,
        5.35,
        3.35,
        [0.65, 1.8, 2.9],
        font_size=7.5,
    )
    page += 1

    failed_rows = summarize_failed_cases(failed_cases)
    slide = prs.slides.add_slide(blank)
    add_title(slide, "24 个未过样例：逐项保留后验验证问题", "CASE LIST", page)
    add_table(slide, failed_rows[:13], 0.45, 1.25, 6.1, 5.65, [1.35, 1.1, 3.65], font_size=6.2)
    add_table(slide, [failed_rows[0], *failed_rows[13:]], 6.75, 1.25, 6.1, 5.65, [1.35, 1.1, 3.65], font_size=6.2)
    page += 1

    slide = prs.slides.add_slide(blank)
    add_title(slide, "如何判断是不是软件逻辑问题", "ROOT-CAUSE METHOD", page)
    add_table(
        slide,
        [
            ["检查层", "判据", "软件动作"],
            ["提资 JSON", "托盘层数、载荷、方钢、厂房、标高是否与报告/源文件一致。", "若解析错，修 intake parser；否则保留差异。"],
            ["谱选择", "SL-1/SL-2、X/Y/Z、阻尼、插值标高是否有 source_ref。", "若谱错，修 selector；若报告用了别的谱，登记冲突。"],
            ["APDL 建模", "三份命令流是否来自标准源族并能解释 LATT/约束/质量。", "若源族选错，修 renderer；历史源冲突不硬改生产源。"],
            ["结果提取", "LIS/OUP/PNG 是否来自正确集合；无全零、UNKNOWN、缺图。", "若节点集错，修 post macro；不得按接近值挑行。"],
            ["评定", "许用值、事故系数、组合比、Excel source_ref 是否明确。", "若公式错，修 evaluator 或 Excel mapping。"],
        ],
        0.65,
        1.35,
        12.05,
        4.4,
        [1.45, 5.25, 5.35],
        font_size=7.8,
    )
    add_card(slide, 0.65, 6.05, 12.05, 0.52, "门禁", "all-zero、UNKNOWN 节点、缺 required 图、未确认公式、真实 ANSYS 失败或 source_map 缺失，都不能发布为正式结论。", fill=RGBColor(254, 243, 242), line=RGBColor(253, 162, 155), body_size=9.5)
    page += 1

    slide = prs.slides.add_slide(blank)
    add_title(slide, "CableTrayAI Job JSON v1：多科室协同输入契约", "JSON CONTRACT", page)
    add_table(
        slide,
        [
            ["JSON 块", "内容", "防错价值"],
            ["job_identity", "提资号/计算批次/后期报告号/来源行", "报告号缺失时仍可计算，后期再绑定。"],
            ["project + spectrum", "项目、厂房、区域、标高、谱文件、阻尼、插值", "不再靠行号或固定厂房标高。"],
            ["support + tray_layers", "支架类型、单双侧、层数、托盘宽度、线荷载、等效密度", "统一提资表达，减少自由文本歧义。"],
            ["materials + section", "材料策略、方钢截面、候选截面选择策略", "截面缺失时按 ratio<1 且接近 1 的规则选择。"],
            ["traceability", "source_ref、创建人、审查状态", "无溯源 JSON 不允许进入 real ANSYS。"],
        ],
        0.75,
        1.35,
        11.85,
        4.2,
        [2.1, 4.8, 4.95],
        font_size=8.0,
    )
    add_card(slide, 0.75, 5.95, 11.85, 0.55, "当前只考虑本地 ANSYS", "JSON 是输入契约，不改变求解后端；后续若接超算，也只替换计算执行层。", fill=RGBColor(236, 253, 243), line=RGBColor(141, 214, 191), body_size=10)
    page += 1

    slide = prs.slides.add_slide(blank)
    add_title(slide, "报告注入边界：模板格式不能变", "REPORT INJECTION", page)
    add_table(
        slide,
        [
            ["对象", "允许替换", "禁止"],
            ["表3-1 反应谱", "表头按厂房/标高改，表格数据来自所选谱文件。", "手填谱值或改表结构。"],
            ["图5.1 / 图5.2", "替换为当前 job 的模型/托臂相关 ANSYS 图片。", "用报告截图或占位图冒充。"],
            ["第六章评定表", "只替换计算值、许用值、比值、结论。", "调整模板排版或自由新增表。"],
            ["附录图片", "来自 figures_manifest.json 中 required 图。", "缺图时生成正式报告。"],
            ["审计文件", "template_report_audit.json 记录所有替换。", "缺 source_ref 仍发布。"],
        ],
        0.75,
        1.4,
        11.85,
        4.2,
        [2.0, 5.0, 4.85],
        font_size=8.2,
    )
    page += 1

    slide = prs.slides.add_slide(blank)
    add_title(slide, "科室审查建议：先审逻辑，再审历史报告", "REVIEW PLAN", page)
    add_table(
        slide,
        [
            ["阶段", "审查材料", "输出"],
            ["1", "三份命令流、JSON v1、result_source_map", "确认建模/计算/提取源逻辑。"],
            ["2", "15 个已审冲突证据", "确认是否列为 legacy conflict 或修订历史报告。"],
            ["3", "24 个未过样例清单", "逐项归因：软件缺陷、历史源差异、报告差异、提资差异。"],
            ["4", "模板注入样例报告", "确认模板格式、表头、图片替换不破坏原格式。"],
            ["5", "新提资试运行", "只看 JSON、ANSYS、result、Excel/公式评定，不再拿旧报告硬套。"],
        ],
        0.8,
        1.45,
        11.7,
        3.75,
        [0.65, 5.6, 5.45],
        font_size=8.2,
    )
    add_text(slide, "结论：初版完全版要优先证明“输入清楚、命令流可审、结果有效、评定有源、报告不改格式”。历史报告 1% 全通过不是唯一目标，尤其在历史源自身冲突时。", 0.85, 5.55, 11.6, 0.55, size=12, bold=True)
    page += 1

    slide = prs.slides.add_slide(blank)
    add_title(slide, "附录：15 个已审冲突记录", "APPENDIX A", page)
    conflict_rows = [["包号", "冲突类型", "证据摘要"]]
    for item in conflicts:
        conflict_rows.append([item.get("package_id", ""), item.get("load_kind") or item.get("evaluation_kind") or "unknown", clean(item.get("reason") or item.get("evidence"), 115)])
    add_table(slide, conflict_rows[:9], 0.45, 1.25, 6.2, 5.65, [1.15, 1.95, 3.1], font_size=5.6)
    add_table(slide, [conflict_rows[0], *conflict_rows[9:]], 6.8, 1.25, 6.05, 5.65, [1.15, 1.95, 2.95], font_size=5.6)
    page += 1

    slide = prs.slides.add_slide(blank)
    add_title(slide, "附录：24 个仍超 1% 后验验证问题", "APPENDIX B", page)
    add_table(slide, failed_rows[:13], 0.45, 1.25, 6.15, 5.65, [1.35, 1.0, 3.8], font_size=6.0)
    add_table(slide, [failed_rows[0], *failed_rows[13:]], 6.8, 1.25, 6.05, 5.65, [1.35, 1.0, 3.7], font_size=6.0)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
